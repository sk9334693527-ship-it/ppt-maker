import os
import json
import random
from datetime import datetime

import pdfplumber
from PIL import Imageimport os
import json
import random
from datetime import datetime

import pdfplumber
from PIL import Image
import pytesseract

import google.generativeai as genai

from telegram import Update, InputFile
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes

from pptx import Presentation
from pptx.util import Pt

# =========================
# 🔐 ADMIN STATES
# =========================
admin_waiting = set()
admin_logged = set()

# =========================
# 🔐 CONFIG
# =========================
BOT_TOKEN = os.getenv("BOT_TOKEN")
ADMIN_ID = int(os.getenv("ADMIN_ID", "0"))

PASSWORD = os.getenv("PASSWORD") or "1234"
PASSWORD = str(PASSWORD).strip()

NUMBER = os.getenv("NUMBER", "Not Set")

GEMINI_API_KEYS = [
    os.getenv("GEMINI_API"),
    os.getenv("GEMINI1_API"),
    os.getenv("GEMINI2_API"),
    os.getenv("GEMINI3_API")
]

DATA_FILE = "data.json"

# =========================
# 💾 DATA
# =========================
def load_data():
    if not os.path.exists(DATA_FILE):
        return {}
    with open(DATA_FILE, "r") as f:
        return json.load(f)

def save_data(data):
    with open(DATA_FILE, "w") as f:
        json.dump(data, f)

data = load_data()

def init_user(user):
    uid = str(user.id)
    if uid not in data:
        data[uid] = {
            "name": user.first_name,
            "credits": 20,
            "history": [],
            "background": False,
            "buffer": []   # 🔥 multi input buffer
        }
        save_data(data)

# =========================
# 🤖 GEMINI
# =========================
def get_model():
    keys = [k for k in GEMINI_API_KEYS if k]
    key = random.choice(keys)
    genai.configure(api_key=key)
    return genai.GenerativeModel("gemini-2.5-flash")

def call_ai(prompt):
    try:
        model = get_model()
        res = model.generate_content(prompt)
        return res.text
    except Exception as e:
        print("AI ERROR:", e)
        return "AI ERROR"

# =========================
# 📄 OCR / PDF
# =========================
def image_to_text(path):
    img = Image.open(path)
    return pytesseract.image_to_string(img)

def pdf_to_text(path):
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    return text

# =========================
# 🧠 PROMPT
# =========================
def build_prompt(text, bilingual=False):
    return f"""
Convert ALL content into MCQ format.

Rules:
- One question per block
- Clean format
- No special symbols

{text}
"""

# =========================
# 🔥 TEXT CLEANER (IMPORTANT FIX)
# =========================
def clean_text(text):
    text = text.replace("\x00", "")
    text = text.replace("\r", "")
    text = text.strip()
    return text[:2000]  # overflow control

# =========================
# 🎞 PPT FIXED
# =========================
def create_ppt(slides, filename):
    prs = Presentation()

    for s in slides:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Question"

        safe_text = clean_text(s)

        tf = content.text_frame
        tf.clear()

        p = tf.add_paragraph()
        p.text = safe_text
        p.font.size = Pt(18)

    prs.save(filename)
    return filename

# =========================
# START
# =========================
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.message.from_user
    init_user(user)
    uid = str(user.id)

    d = data[uid]

    await update.message.reply_text(f"""
👤 {d['name']}
🆔 {uid}
💳 {d['credits']}

/objective
/objective2
/admin

📞 {NUMBER}
""")

# =========================
# COMMANDS
# =========================
async def objective(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["mode"] = "objective"
    await update.message.reply_text("Send multiple text, then send /done")

async def done(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = str(update.message.from_user.id)

    if not data[uid]["buffer"]:
        await update.message.reply_text("No data")
        return

    full_text = "\n".join(data[uid]["buffer"])
    data[uid]["buffer"] = []

    prompt = build_prompt(full_text)
    ai_text = call_ai(prompt)

    slides = [s.strip() for s in ai_text.split("\n\n") if s.strip()]

    ppt = f"{uid}.pptx"
    create_ppt(slides, ppt)

    await update.message.reply_document(InputFile(ppt, filename="output.pptx"))

# =========================
# HANDLE
# =========================
async def handle(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.message.from_user
    uid = str(user.id)
    init_user(user)

    if update.message.text:
        data[uid]["buffer"].append(update.message.text)
        save_data(data)
        await update.message.reply_text("Added")

# =========================
# MAIN
# =========================
def main():
    app = Application.builder().token(BOT_TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("objective", objective))
    app.add_handler(CommandHandler("done", done))

    app.add_handler(MessageHandler(filters.ALL, handle))

    print("Bot running...")
    app.run_polling()

if __name__ == "__main__":
    main()
import pytesseract

import google.generativeai as genai

from telegram import Update, InputFile
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes

from pptx import Presentation
from pptx.util import Pt

# =========================
# 🔐 ADMIN STATES
# =========================
admin_waiting = set()
admin_logged = set()

# =========================
# 🔐 CONFIG
# =========================
BOT_TOKEN = os.getenv("BOT_TOKEN")
ADMIN_ID = int(os.getenv("ADMIN_ID", "0"))

PASSWORD = os.getenv("PASSWORD")
if PASSWORD is None:
    PASSWORD = "1234"
PASSWORD = str(PASSWORD).strip()

NUMBER = os.getenv("NUMBER", "Not Set")

GEMINI_API_KEYS = [
    os.getenv("GEMINI_API"),
    os.getenv("GEMINI1_API"),
    os.getenv("GEMINI2_API"),
    os.getenv("GEMINI3_API")
]

DATA_FILE = "data.json"

# =========================
# 💾 DATA SYSTEM
# =========================
def load_data():
    if not os.path.exists(DATA_FILE):
        return {}
    with open(DATA_FILE, "r") as f:
        return json.load(f)

def save_data(data):
    with open(DATA_FILE, "w") as f:
        json.dump(data, f)

data = load_data()

def init_user(user):
    uid = str(user.id)
    if uid not in data:
        data[uid] = {
            "name": user.first_name,
            "credits": 20,
            "history": [],
            "background": False
        }
        save_data(data)

# =========================
# 🤖 GEMINI AI
# =========================
def get_model():
    keys = [k for k in GEMINI_API_KEYS if k]
    if not keys:
        raise Exception("No Gemini API found")

    key = random.choice(keys)
    genai.configure(api_key=key)
    return genai.GenerativeModel("gemini-2.5-flash")

def call_ai(prompt):
    try:
        model = get_model()
        res = model.generate_content(prompt)
        return res.text
    except Exception as e:
        print("AI ERROR:", e)
        return "AI ERROR"

# =========================
# 📄 OCR
# =========================
def image_to_text(path):
    img = Image.open(path)
    return pytesseract.image_to_string(img)

# =========================
# 📄 PDF
# =========================
def pdf_to_text(path):
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    return text

# =========================
# 🧠 PROMPT
# =========================
def build_prompt(text, bilingual=False):
    if bilingual:
        return f"""
Convert into MCQ PPT format.

Hindi + English
One question per slide

{text}
"""
    else:
        return f"""
Convert into MCQ PPT format.

English only
One question per slide

{text}
"""

# =========================
# 🎞 PPT
# =========================
def create_ppt(slides, filename):
    prs = Presentation()

    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Question"
        slide.placeholders[1].text = s

        for p in slide.placeholders[1].text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(18)

    prs.save(filename)
    return filename

# =========================
# /START
# =========================
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.message.from_user
    init_user(user)
    uid = str(user.id)

    d = data[uid]

    await update.message.reply_text(f"""
👤 Name: {d['name']}
🆔 ID: {uid}
💳 Credits: {d['credits']}

📌 MENU:
/objective
/objective2
/make
/background
/dbackground
/admin

📞 {NUMBER}
""")

# =========================
# COMMANDS (UPDATED)
# =========================
async def objective(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["mode"] = "objective"
    context.user_data["buffer"] = []

    await update.message.reply_text("📥 Send multiple texts\nThen send /make")

async def objective2(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["mode"] = "objective2"
    context.user_data["buffer"] = []

    await update.message.reply_text("📥 Send Hindi+English texts\nThen /make")

async def make(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.message.from_user
    uid = str(user.id)

    if "buffer" not in context.user_data or not context.user_data["buffer"]:
        await update.message.reply_text("❌ No data found")
        return

    full_text = "\n".join(context.user_data["buffer"])

    mode = context.user_data.get("mode", "objective")
    prompt = build_prompt(full_text, bilingual=(mode == "objective2"))

    ai_text = call_ai(prompt)

    slides = [s.strip() for s in ai_text.split("\n\n") if s.strip()]

    ppt = f"{uid}_{int(datetime.now().timestamp())}.pptx"
    create_ppt(slides, ppt)

    data[uid]["credits"] -= len(slides)
    data[uid]["history"].append({
        "time": str(datetime.now()),
        "slides": len(slides)
    })
    save_data(data)

    context.user_data["buffer"] = []

    await update.message.reply_document(
        document=InputFile(ppt, filename=ppt)
    )

async def background(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text("Send background ppt")

async def dbackground(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = str(update.message.from_user.id)
    data[uid]["background"] = False
    save_data(data)
    await update.message.reply_text("Background removed")

async def admin(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.message.from_user.id
    admin_waiting.add(uid)
    await update.message.reply_text("🔐 Send password")

# =========================
# MAIN HANDLE (UPDATED)
# =========================
async def handle(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.message.from_user
    uid = user.id
    init_user(user)

    text_msg = update.message.text.strip() if update.message.text else ""

    # ADMIN LOGIN
    if uid in admin_waiting:
        if text_msg == PASSWORD and uid == ADMIN_ID:
            admin_waiting.remove(uid)
            admin_logged.add(uid)
            await update.message.reply_text("✅ ADMIN PANEL")
        else:
            admin_waiting.remove(uid)
            await update.message.reply_text("❌ Wrong password")
        return

    # 🔥 MULTI MESSAGE COLLECT
    if update.message.text and context.user_data.get("mode") in ["objective", "objective2"]:
        if update.message.text != "/make":
            context.user_data.setdefault("buffer", []).append(update.message.text)
            await update.message.reply_text("✅ Added")
            return

    # CREDIT CHECK
    uid_str = str(uid)
    if data[uid_str]["credits"] <= 0:
        await update.message.reply_text("No credits")
        return

# =========================
# MAIN
# =========================
def main():
    app = Application.builder().token(BOT_TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("objective", objective))
    app.add_handler(CommandHandler("objective2", objective2))
    app.add_handler(CommandHandler("make", make))
    app.add_handler(CommandHandler("background", background))
    app.add_handler(CommandHandler("dbackground", dbackground))
    app.add_handler(CommandHandler("admin", admin))

    app.add_handler(MessageHandler(filters.ALL, handle))

    print("Bot running...")
    app.run_polling()

if __name__ == "__main__":
    main()
